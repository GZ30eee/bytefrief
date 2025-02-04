import streamlit as st
import requests
import spacy
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.decomposition import TruncatedSVD
import docx
import PyPDF2
from pptx import Presentation
from io import BytesIO
from spacy.cli import download

# Try to load the model; if not found, download it
try:
    nlp = spacy.load("en_core_web_sm")
except OSError:
    print("Downloading 'en_core_web_sm' model...")
    download("en_core_web_sm")  # Download the model
    nlp = spacy.load("en_core_web_sm")  # Load the model after download


# Function to fetch random text
def fetch_random_paragraph():
    try:
        response = requests.get("https://baconipsum.com/api/?type=all-meat&paras=1&start-with-lorem=1")
        if response.status_code == 200:
            return ' '.join(response.json())
        return "Could not fetch content. Please try again."
    except requests.exceptions.RequestException:
        return "Connection error. Please check your internet connection."

# Custom summarizer function
def custom_summarizer(text, target_word_count, keywords=None):
    doc = nlp(text)
    sentences = [sent.text.strip() for sent in doc.sents]

    if len(sentences) < 2:
        return text, len(sentences), len(text.split())

    vectorizer = TfidfVectorizer(stop_words='english', ngram_range=(1, 2))
    X = vectorizer.fit_transform(sentences)

    svd = TruncatedSVD(n_components=1, random_state=42)
    svd.fit(X)
    scores = svd.components_[0]

    ranked_sentences = []
    for i, sent in enumerate(sentences):
        score = scores[i] * 0.7 + (1 / (i + 1)) * 0.3
        ranked_sentences.append((score, sent))

    ranked_sentences.sort(reverse=True)

    summary = []
    current_length = 0
    for score, sent in ranked_sentences:
        words = len(sent.split())
        if current_length + words <= target_word_count:
            summary.append(sent)
            current_length += words
        if current_length >= target_word_count:
            break

    return ' '.join(summary), len(summary), current_length

# Function to convert summary to bullet points
def convert_to_bullets(summary_text):
    doc = nlp(summary_text)
    return '\n\n'.join([f"• {sent.text}" for sent in doc.sents])

# Extract text from uploaded files
def extract_text_from_uploaded_file(uploaded_file):
    file_extension = uploaded_file.name.split('.')[-1].lower()
    text = ""

    if file_extension == 'txt':
        text = uploaded_file.read().decode("utf-8")
    elif file_extension == 'pdf':
        reader = PyPDF2.PdfReader(uploaded_file)
        text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
    elif file_extension == 'docx':
        doc = docx.Document(uploaded_file)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif file_extension in ['ppt', 'pptx']:
        prs = Presentation(uploaded_file)
        text = "\n".join([slide.shapes.title.text if slide.shapes.title else '' for slide in prs.slides])

    return text

# Function to export summary as a DOCX file
def export_to_docx(summary_text):
    doc = docx.Document()
    doc.add_paragraph(summary_text)
    doc_io = BytesIO()
    doc.save(doc_io)
    doc_io.seek(0)
    return doc_io

# Function to process custom commands
def process_custom_command(command, text):
    doc = nlp(text)

    if command.lower() == "give a title":
        return list(doc.sents)[0].text.strip() if len(text) > 0 else "Untitled"
    elif command.lower() == "generate a conclusion":
        return "In conclusion, " + " ".join([sent.text for sent in doc.sents if len(sent.text) > 20])[:200]
    elif command.lower() == "make it academic":
        return " ".join([token.text.upper() if token.pos_ == "NOUN" else token.text for token in doc])
    else:
        return "Command not recognized. Please try again."

# Streamlit UI
st.title("Smart Text Summarizer")

# Select output mode
mode = st.selectbox("Select Output Format", ["Paragraph", "Bullet Points", "Custom Command"])

# File uploader
uploaded_file = st.file_uploader("Upload text file", type=["txt", "ppt", "pptx", "docx", "pdf"])

if "text" not in st.session_state:
    st.session_state.text = ""

if uploaded_file:
    st.session_state.text = extract_text_from_uploaded_file(uploaded_file)
else:
    text_area = st.text_area("Paste your text here", height=200, value=st.session_state.text)
    st.session_state.text = text_area

# Random text button
if st.button("Get Random Text"):
    st.session_state.text = fetch_random_paragraph()

# Custom command selection
if mode == "Custom Command":
    custom_cmd = st.selectbox("Choose a Command", ["Give a Title", "Generate a Conclusion", "Make it Academic"])

# Keyword selection
keywords = []
if mode == "Paragraph" or mode == "Custom Command":
    all_keywords = [token.text for token in nlp(st.session_state.text) if not token.is_stop and token.pos_ in ["NOUN", "PROPN"]]
    keywords = st.multiselect("Select keywords to enhance the summary", all_keywords)

# **Slider for word count**
target_words = st.slider("Target summary length", 50, 300, 150, help="Number of words for the summary")

# Generate summary button
if st.button("Generate Summary"):
    if st.session_state.text:
        if mode == "Paragraph":
            result, num_sentences, num_words = custom_summarizer(st.session_state.text, target_words, keywords)
        elif mode == "Bullet Points":
            summary, num_sentences, num_words = custom_summarizer(st.session_state.text, target_words)
            result = convert_to_bullets(summary)
        elif mode == "Custom Command":
            result = process_custom_command(custom_cmd, st.session_state.text) if custom_cmd else "Please enter a command"
            num_sentences = len(result.split('.'))
            num_words = len(result.split())

        st.subheader(f"Result - {num_sentences} Sentences, {num_words} Words")
        st.write(result if result else "No output generated")

        # Statistics
        with st.expander("Statistics"):
            word_count = len(st.session_state.text.split())
            characters = len(st.session_state.text)
            reduction = round((1 - (num_words / word_count)) * 100, 2) if word_count > 0 else 0
            st.write(f"Word count: {word_count}")
            st.write(f"Sentence count: {num_sentences}")
            st.write(f"Characters: {characters}")
            st.write(f"Reduction: {reduction}%")

        
# Clear button
if st.button("Clear All"):
    st.session_state.text = ""
    st.experimental_rerun()
